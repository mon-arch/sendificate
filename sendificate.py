import glob
import smtplib
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from tkinter import filedialog, scrolledtext

import comtypes.client
import pandas as pd
from pptx import Presentation

import os
from PIL import ImageTk, Image

import tkinter as tk
from ttkbootstrap import Style


def exit_application():
    root.destroy()
def check_credentials():
    entered_username = username_entry.get()
    entered_password = password_entry.get()

    # Kullanıcı adı ve şifre kontrolü
    if entered_username == "emusoft.ai" and entered_password == "18811938":
        root.destroy()
        Main_Application()

    else:
        access_granted_label.config(text="Giriş başarısız! Lütfen tekrar deneyin.", fg="red")
        username_entry.delete(0, tk.END)
        password_entry.delete(0, tk.END)

def open_contact_window():
    contact_window = tk.Toplevel(root)
    contact_window.title("İletişim Bilgileri")

    developer_info = """
    Geliştiriciler:
    - Kaan Acar
    - Mehmet Oktar Önder
    - Elifnaz Çelik

    İletişim:
    Email: kaanacar@skiff.com
    Email: oktar.onder@gmail.com
    Email: d.elifcelik@gmail.com
    """

    developer_info_label = tk.Label(contact_window, text=developer_info, font=('Garamond', 14, "bold"))
    developer_info_label.pack()
    contact_window.iconbitmap("logo.ico")
    contact_window.minsize(300, 300)

def center_window(window, width, height):
    # Ekranın genişliği ve yüksekliği
    screen_width = window.winfo_screenwidth()
    screen_height = window.winfo_screenheight()

    # Pencereyi ortalamak için hesaplamalar
    x = (screen_width // 2) - (width // 2)
    y = (screen_height // 2) - (height // 2)

    window.geometry('{}x{}+{}+{}'.format(width, height, x, y))

window_width = 450
window_height = 500

# Ana pencere oluşturma
root = tk.Tk()
root.overrideredirect(1)
window_width = 450
window_height = 500
root.title("emusoft.ai/Yazılım ve Yapay Zeka Geliştirme Kulübü")
root.iconbitmap("logo.ico")


# Stili Tkinter için ayarla
style = Style(theme='flatly')

# Geliştirici bilgisi resmi
developer_info_img = Image.open("emusoft.ai.png")
developer_info_img = developer_info_img.resize((200, 200), Image.LANCZOS)
developer_info_img = ImageTk.PhotoImage(developer_info_img)

# Geliştirici bilgisi paneli
developer_info_panel = tk.Label(root, image=developer_info_img)
developer_info_panel.pack()

# Geliştirici yazısı
developer_text = tk.Label(root, text="Sendificate", font=('Brush Script MT', 50, "bold"), fg='dark blue')
developer_text.pack()

# Kullanıcı adı girişi
username_label = tk.Label(root, text="Kullanıcı Adı:", font=("Now", 10, "bold"))
username_label.pack()

username_entry = tk.Entry(root)
username_entry.pack()

# Şifre girişi
password_label = tk.Label(root, text="Şifre:", font=("Now", 10, "bold"))
password_label.pack()

password_entry = tk.Entry(root, show="*")
password_entry.pack()

# Ek Yazı
additional_text = tk.Label(root, text="", font=('Garamond', 5, "bold"))
additional_text.pack()

# Giriş butonu
submit_button = tk.Button(root, text="Giriş Yap", command=check_credentials)
submit_button.pack()

# Bilgilendirme etiketi
access_granted_label = tk.Label(root, text="", fg="black")
access_granted_label.pack()

# Ek metin ve İletişim butonu
additional_text = tk.Label(root, text="Soru ve iş birlikleriniz için; ", font=('Garamond', 14, "bold"))
additional_text.pack()

contact_button = tk.Button(root, text="İletişim", command=open_contact_window, font=('Garamond', 12, "bold"))
contact_button.pack()

additional_text = tk.Label(root, text="", font=('Garamond', 14, "bold"))
additional_text.pack()

exit_button = tk.Button(root, text="X", command=exit_application, bg='red', fg='white', font=('Arial', 10, 'bold'))
exit_button.place(relx=1, x=-2, y=2, anchor='ne')

center_window(root, window_width, window_height)

def Main_Application():

    def print_to_console(message):
        console.config(state=tk.NORMAL)
        console.insert(tk.END, message + "\n")
        console.config(state=tk.DISABLED)
        console.yview(tk.END)

    def create_powerpoint(excel_path, template_path, output_folder):
        excel_data = pd.read_excel(excel_path)
        template = Presentation(template_path)
        column_name = 'AD'  # Değiştirmeniz gereken sütun adı
        column_data = excel_data[column_name]
        target_text = '<<FULL NAME>>'
        os.makedirs(output_folder, exist_ok=True)

        for index, value in enumerate(column_data):
            # Türkçe karakterlerin dönüşümü
            turkish_to_english = str.maketrans("iı", "İI")

            for i, char in enumerate(value):
                if char.lower() in "iı":
                    value = value[:i] + char.translate(turkish_to_english) + value[i + 1:]
                    
            value_upper = str(value).upper()
            presentation = Presentation(template_path)

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if target_text in run.text:
                                    run.text = run.text.replace(target_text, value_upper)

            # Yeni dosya isimleri için formatlanmış isim oluşturma
            file_number = str(index + 1).zfill(4)  # Dört basamaklı numara oluşturma
            pptx_file_name = f"{file_number}.pptx"

            output_path = os.path.join(output_folder, pptx_file_name)
            print_to_console(f"'{pptx_file_name}' adlı PowerPoint sunusu '{output_folder}' klasörüne kaydedildi.")
            presentation.save(output_path)
        print_to_console(f"Tüm PowerPoint sunuları '{output_folder}' klasörüne kaydedildi.")

    def select_excel():
        file_path = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx")])
        if file_path:
            excel_entry.delete(0, tk.END)
            excel_entry.insert(tk.END, file_path)

    def select_template():
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            template_entry.delete(0, tk.END)
            template_entry.insert(tk.END, file_path)

    def select_output_folder():
        folder_path = filedialog.askdirectory()
        if folder_path:
            output_entry.delete(0, tk.END)
            output_entry.insert(tk.END, folder_path)

    def generate_powerpoints():

        excel_path = excel_entry.get()
        template_path = template_entry.get()
        output_folder = output_entry.get()

        if all([excel_path, template_path, output_folder]):
            create_powerpoint(excel_path, template_path, output_folder)

    def convert_pptx_to_pdf(input_file_path, output_folder_path):
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = True

        pptx_file = os.path.abspath(input_file_path)
        pdf_folder = os.path.abspath(output_folder_path)

        presentation = powerpoint.Presentations.Open(pptx_file)
        pdf_file = os.path.join(pdf_folder, os.path.splitext(os.path.basename(pptx_file))[0] + ".pdf")
        presentation.SaveAs(pdf_file, 32)
        presentation.Close()

        powerpoint.Quit()
        print_to_console(f"PDF'in kaydedildiği konum: {pdf_file}")

    def convert_folder_to_pdf(folder_path, output_folder_path):
        for file in os.listdir(folder_path):
            if file.endswith(".pptx"):
                convert_pptx_to_pdf(os.path.join(folder_path, file), output_folder_path)
                print_to_console(f"'{file}' adlı dosya dönüştürüldü.")
        print_to_console("Tüm dönüştürme işlemleri tamamlandı.")

    def select_input_folder():
        folder_path = filedialog.askdirectory()
        if folder_path:
            folder_path_entry.delete(0, tk.END)
            folder_path_entry.insert(tk.END, folder_path)

    def select_output_folder_for_pdf():
        folder_path = filedialog.askdirectory()
        if folder_path:
            output_folder_entry.delete(0, tk.END)
            output_folder_entry.insert(tk.END, folder_path)

    def convert():
        input_folder_path = folder_path_entry.get()
        output_folder_path = output_folder_entry.get()

        if input_folder_path and output_folder_path:
            convert_folder_to_pdf(input_folder_path, output_folder_path)

    def send_email_with_pdf_and_message(pdf_file, recipient_email, sender_email, sender_password, subject,
                                        message_content):
        # SMTP bağlantısı ve oturum açma
        smtp_server = 'smtp-mail.outlook.com'  # Outlook SMTP sunucu adresi
        smtp_port = 587  # SMTP port numarası

        # PDF dosyasını oku
        with open(pdf_file, 'rb') as file:
            pdf_content = file.read()

        # E-posta oluşturma
        message = MIMEMultipart()
        message['From'] = sender_email
        message['To'] = recipient_email
        message['Subject'] = subject

        # Metin içeriğini ekle
        message_text = MIMEText(message_content)
        message.attach(message_text)

        # PDF dosyasını ekle
        pdf_attachment = MIMEApplication(pdf_content)
        pdf_attachment.add_header('Content-Disposition', 'attachment', filename=os.path.basename(pdf_file))
        message.attach(pdf_attachment)

        # E-posta gönderme işlemi
        try:
            server = smtplib.SMTP(smtp_server, smtp_port)
            server.starttls()
            server.login(sender_email, sender_password)
            server.sendmail(sender_email, recipient_email, message.as_string())
            print_to_console(f"PDF dosyası {recipient_email} adresine gönderildi.")
        except Exception as e:
            print_to_console(f"Hata oluştu: {str(e)}")
        finally:
            server.quit()

    def select_excel_for_email():
        file_path = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx")])
        if file_path:
            excel_entry.delete(0, tk.END)
            excel_entry.insert(tk.END, file_path)

    def select_pdf_folder_for_email():
        folder_path = filedialog.askdirectory()
        if folder_path:
            pdf_folder_entry.delete(0, tk.END)
            pdf_folder_entry.insert(tk.END, folder_path)

    def send_emails():
        excel_file = excel_entry.get()
        pdf_folder = pdf_folder_entry.get()
        sender_email = sender_email_entry.get()
        sender_password = sender_password_entry.get()
        subject = subject_entry.get()
        message_content = message_content_text.get("1.0", tk.END)

        df = pd.read_excel(excel_file)
        pdf_files = glob.glob(f"{pdf_folder}/*.pdf")

        for index, row in df.iterrows():
            email = row['Email']
            pdf_index = index % len(pdf_files)
            pdf_to_send = pdf_files[pdf_index]

            send_email_with_pdf_and_message(pdf_to_send, email, sender_email, sender_password, subject, message_content)

    root = tk.Tk()
    root.title("emusoft.ai/Yazılım ve Yapay Zeka Geliştirme Kulübü")

    root.resizable(False, False)

    # Logo eklemek için ikon dosyasının yolu
    logo_path = "logo.ico"  # Örnek bir .ico dosyası yolu

    # İkonu pencereye ekle
    root.iconbitmap(logo_path)

    # PowerPoint Oluşturma Bölümü
    frame1 = tk.LabelFrame(root, text="PowerPoint Oluşturucu", padx=10, pady=10)
    frame1.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")

    excel_label = tk.Label(frame1, text="Excel Dosyasını Seçin:")
    excel_label.grid(row=0, column=0, sticky="w")

    excel_entry = tk.Entry(frame1, width=50)
    excel_entry.grid(row=0, column=1)

    excel_button = tk.Button(frame1, text="Seç", command=select_excel)
    excel_button.grid(row=0, column=2)

    template_label = tk.Label(frame1, text="PowerPoint Şablonunu Seçin:")
    template_label.grid(row=1, column=0, sticky="w")

    template_entry = tk.Entry(frame1, width=50)
    template_entry.grid(row=1, column=1)

    template_button = tk.Button(frame1, text="Seç", command=select_template)
    template_button.grid(row=1, column=2)

    output_label = tk.Label(frame1, text="Powerpointlerin Kaydedileceği Klasörü Seçin:")
    output_label.grid(row=2, column=0, sticky="w")

    output_entry = tk.Entry(frame1, width=50)
    output_entry.grid(row=2, column=1)

    output_button = tk.Button(frame1, text="Seç", command=select_output_folder)
    output_button.grid(row=2, column=2)

    generate_button = tk.Button(frame1, text="PowerPoint'leri Oluştur", command=generate_powerpoints)
    generate_button.grid(row=3, columnspan=3)

    # PowerPoint -> PDF Dönüştürme Bölümü
    frame2 = tk.LabelFrame(root, text="PowerPoint -> Sertifika Dönüştürücü", padx=10, pady=10)
    frame2.grid(row=0, column=1, padx=10, pady=10, sticky="nsew")

    input_folder_label = tk.Label(frame2, text="PowerPoint Klasörünü Seçin:")
    input_folder_label.grid(row=0, column=0, sticky="w")

    folder_path_entry = tk.Entry(frame2, width=50)
    folder_path_entry.grid(row=0, column=1)

    input_folder_button = tk.Button(frame2, text="Seç", command=select_input_folder)
    input_folder_button.grid(row=0, column=2)

    output_folder_label = tk.Label(frame2, text="Sertifika Klasörünü Seçin:")
    output_folder_label.grid(row=1, column=0, sticky="w")

    output_folder_entry = tk.Entry(frame2, width=50)
    output_folder_entry.grid(row=1, column=1)

    output_folder_button = tk.Button(frame2, text="Seç", command=select_output_folder_for_pdf)
    output_folder_button.grid(row=1, column=2)

    convert_button = tk.Button(frame2, text="PowerPoint'ten Sertifikaya Dönüştür", command=convert)
    convert_button.grid(row=2, columnspan=3)

    # Email Gönderme Bölümü
    frame3 = tk.LabelFrame(root, text="Email Gönderme", padx=10, pady=10)
    frame3.grid(row=1, column=0, padx=10, pady=10, sticky="nsew")

    # Sender Email
    sender_email_label = tk.Label(frame3, text="Email:")
    sender_email_label.grid(row=0, column=0, sticky="w")
    sender_email_entry = tk.Entry(frame3, width=50)
    sender_email_entry.grid(row=0, column=1)

    # Sender Password
    sender_password_label = tk.Label(frame3, text="Şifre:")
    sender_password_label.grid(row=1, column=0, sticky="w")
    sender_password_entry = tk.Entry(frame3, show="*", width=50)
    sender_password_entry.grid(row=1, column=1)

    # Subject
    subject_label = tk.Label(frame3, text="Konu:")
    subject_label.grid(row=2, column=0, sticky="w")
    subject_entry = tk.Entry(frame3, width=50)
    subject_entry.grid(row=2, column=1)

    # Message Content
    message_content_label = tk.Label(frame3, text="Mesaj İçeriği:")
    message_content_label.grid(row=3, column=0, sticky="w")
    message_content_text = tk.Text(frame3, width=50, height=10)
    message_content_text.grid(row=3, column=1)

    # Excel File
    excel_label_email = tk.Label(frame3, text="Excel Dosyası:")
    excel_label_email.grid(row=4, column=0, sticky="w")
    excel_entry = tk.Entry(frame3, width=50)
    excel_entry.grid(row=4, column=1)
    excel_button_email = tk.Button(frame3, text="Seç", command=select_excel_for_email)
    excel_button_email.grid(row=4, column=2)

    # PDF Folder
    pdf_folder_label_email = tk.Label(frame3, text="Sertifika Klasörü:")
    pdf_folder_label_email.grid(row=5, column=0, sticky="w")
    pdf_folder_entry = tk.Entry(frame3, width=50)
    pdf_folder_entry.grid(row=5, column=1)
    pdf_folder_button_email = tk.Button(frame3, text="Seç", command=select_pdf_folder_for_email)
    pdf_folder_button_email.grid(row=5, column=2)

    # Send Button
    send_button = tk.Button(frame3, text="Emailleri Gönder", command=send_emails)
    send_button.grid(row=6, columnspan=3)

    frame4 = tk.LabelFrame(root, text="Konsol", padx=10, pady=10)
    frame4.grid(row=1, column=1, padx=10, pady=10, sticky="nsew")

    console = scrolledtext.ScrolledText(frame4, wrap=tk.WORD, width=80, height=15)
    console.pack()

    # Örnek bir kullanım
    print_to_console(
        "/*\nGeliştiriciler\n---------------\nKaan ACAR - kaanacar@skiff.com\nMehmet Oktar ÖNDER - oktar.onder@gmail.com\nElifnaz ÇELİK - d.elifcelik@gmail.com\n*/")

    root.mainloop()

root.mainloop()
